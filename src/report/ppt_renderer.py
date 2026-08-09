from __future__ import annotations

import os
from dataclasses import dataclass, field
from datetime import date, datetime, timezone
from io import BytesIO
from pathlib import Path
from typing import Optional
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt

from .models import ReportCitationDraft, ReportSectionDraft, ReportSectionStatus
from .pdf_renderer import normalize_pdf_text


DEFAULT_PPT_TITLE = "\uc77c\uc77c \uc0b0\uc5c5 \ub3d9\ud5a5 \ubcf4\uace0\uc11c"
DEFAULT_PPT_FONT_NAME = "Malgun Gothic"
PROJECT_ROOT = Path(__file__).resolve().parents[2]
DEFAULT_PPT_LOGO_PATH = PROJECT_ROOT / "assets" / "mySUNI.png"
PPT_LOGO_ENV = "MYWIKI_PPT_LOGO_PATH"
BRAND_NAVY = RGBColor(27, 43, 65)
BRAND_ORANGE = RGBColor(242, 101, 34)
BRAND_GOLD = RGBColor(247, 181, 0)
BRAND_SAND = RGBColor(247, 243, 236)
BRAND_MIST = RGBColor(239, 244, 247)
BRAND_TEXT = RGBColor(52, 61, 72)
BRAND_MUTED = RGBColor(104, 117, 132)

SAFE_LEFT = 0.55
SAFE_RIGHT = 0.55
SAFE_TOP = 0.45
SAFE_BOTTOM = 0.45
SAFE_WIDTH_4_3 = 10.0 - SAFE_LEFT - SAFE_RIGHT
TITLE_MAX_HEIGHT = 0.8
TITLE_DIVIDER_GAP = 0.2
MAX_SUMMARY_CHARS = 280
MAX_BULLET_CHARS = 45
MAX_COLUMN_ITEMS = 3


@dataclass(frozen=True)
class PptLayout:
    title_font_size: int = 34
    slide_title_font_size: int = 28
    card_title_font_size: int = 19
    body_font_size: float = 16
    meta_font_size: int = 12
    max_evidences_per_section: int = 6
    max_evidences_per_slide: int = 3
    max_agenda_items_per_slide: int = 5
    slide_width_inches: float = 10.0
    slide_height_inches: float = 7.5


@dataclass(frozen=True)
class PptEvidenceLine:
    document_version_id: str
    citation_order: int
    evidence_text: str
    relevance_score: Optional[float] = None
    document_title: Optional[str] = None
    source_name: Optional[str] = None
    published_at: Optional[str] = None
    source_url: Optional[str] = None


@dataclass(frozen=True)
class PptSection:
    category: str
    title: str
    status: str
    importance_score: Optional[int]
    impact_direction: Optional[str]
    time_horizon: Optional[str]
    current_summary: Optional[str]
    key_facts: tuple[str, ...] = ()
    historical_context: tuple[str, ...] = ()
    implications: tuple[str, ...] = ()
    watch_points: tuple[str, ...] = ()
    evidences: tuple[PptEvidenceLine, ...] = ()


@dataclass(frozen=True)
class PptReportDocument:
    title: str
    subtitle: str
    generated_at: str
    report_date: str
    version: int
    layout: PptLayout = field(default_factory=PptLayout)
    sections: tuple[PptSection, ...] = ()


DEFAULT_DAILY_REPORT_PPT_LAYOUT = PptLayout()


def build_daily_report_ppt_document(
    *,
    report_key: str,
    version: int,
    sections: list[ReportSectionDraft],
    generated_at: Optional[str] = None,
    report_date: date | str | None = None,
    title: Optional[str] = None,
    layout: PptLayout = DEFAULT_DAILY_REPORT_PPT_LAYOUT,
) -> PptReportDocument:
    completed_sections = [
        _to_ppt_section(section, max_evidences=layout.max_evidences_per_section)
        for section in sections
        if section.status == ReportSectionStatus.COMPLETED
    ]
    normalized_generated_at = normalize_pdf_text(generated_at or _utc_now_iso())
    normalized_report_date = normalize_pdf_text(
        _resolve_report_date_text(report_date=report_date, generated_at=normalized_generated_at)
    )
    return PptReportDocument(
        title=normalize_pdf_text(title or DEFAULT_PPT_TITLE),
        subtitle=normalize_pdf_text(report_key),
        generated_at=normalized_generated_at,
        report_date=normalized_report_date,
        version=version,
        layout=layout,
        sections=tuple(completed_sections),
    )


def build_daily_report_ppt_filename(*, report_key: str, version: int) -> str:
    normalized_key = normalize_pdf_text(report_key).strip().replace(" ", "-")
    return f"{normalized_key}-v{version}.pptx"


def render_daily_report_ppt(document: PptReportDocument) -> bytes:
    normalized_document = _normalize_document(document)
    presentation = Presentation()
    presentation.slide_width = Inches(normalized_document.layout.slide_width_inches)
    presentation.slide_height = Inches(normalized_document.layout.slide_height_inches)

    _add_cover_slide(presentation, normalized_document)
    _add_agenda_slides(presentation, normalized_document)

    if not normalized_document.sections:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        _fill_background(slide, BRAND_MIST)
        _add_slide_title(slide, "\uc644\ub8cc\ub41c \uc139\uc158", normalized_document.layout)
        _add_text_block(
            slide,
            left=0.9,
            top=1.8,
            width=6.5,
            height=1.2,
            text="\uc644\ub8cc\ub41c \uc139\uc158\uc774 \uc5c6\uc2b5\ub2c8\ub2e4.",
            font_size=20,
            layout=normalized_document.layout,
            color=BRAND_TEXT,
        )
    else:
        for index, section in enumerate(normalized_document.sections, start=1):
            _add_section_overview_slide(presentation, normalized_document, index, section)
            _add_section_highlight_slide(presentation, normalized_document, index, section)
        all_source_items = _collect_source_items(normalized_document.sections)
        if all_source_items:
            for chunk_index, source_chunk in enumerate(_chunked_text(all_source_items, 5), start=1):
                _add_sources_slide(presentation, normalized_document, chunk_index, source_chunk)

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _utc_now_iso() -> str:
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat()


def _resolve_report_date_text(*, report_date: date | str | None, generated_at: str) -> str:
    if isinstance(report_date, date):
        return report_date.isoformat()
    if isinstance(report_date, str) and report_date.strip():
        raw = report_date.strip()
        if "T" in raw:
            try:
                return datetime.fromisoformat(raw.replace("Z", "+00:00")).date().isoformat()
            except ValueError:
                return raw.split("T", 1)[0]
        return raw
    try:
        return datetime.fromisoformat(generated_at.replace("Z", "+00:00")).date().isoformat()
    except ValueError:
        return generated_at.split("T", 1)[0]


def _normalize_document(document: PptReportDocument) -> PptReportDocument:
    return PptReportDocument(
        title=normalize_pdf_text(document.title),
        subtitle=normalize_pdf_text(document.subtitle),
        generated_at=normalize_pdf_text(document.generated_at),
        report_date=normalize_pdf_text(document.report_date),
        version=document.version,
        layout=document.layout,
        sections=tuple(
            PptSection(
                category=normalize_pdf_text(section.category),
                title=normalize_pdf_text(section.title),
                status=normalize_pdf_text(section.status),
                importance_score=section.importance_score,
                impact_direction=normalize_pdf_text(section.impact_direction) if section.impact_direction else None,
                time_horizon=normalize_pdf_text(section.time_horizon) if section.time_horizon else None,
                current_summary=normalize_pdf_text(section.current_summary) if section.current_summary else None,
                key_facts=tuple(normalize_pdf_text(item) for item in section.key_facts),
                historical_context=tuple(normalize_pdf_text(item) for item in section.historical_context),
                implications=tuple(normalize_pdf_text(item) for item in section.implications),
                watch_points=tuple(normalize_pdf_text(item) for item in section.watch_points),
                evidences=tuple(
                    PptEvidenceLine(
                        document_version_id=normalize_pdf_text(evidence.document_version_id),
                        citation_order=evidence.citation_order,
                        evidence_text=normalize_pdf_text(evidence.evidence_text or ""),
                        relevance_score=evidence.relevance_score,
                        document_title=normalize_pdf_text(evidence.document_title) if evidence.document_title else None,
                        source_name=normalize_pdf_text(evidence.source_name) if evidence.source_name else None,
                        published_at=normalize_pdf_text(evidence.published_at) if evidence.published_at else None,
                        source_url=normalize_pdf_text(evidence.source_url) if evidence.source_url else None,
                    )
                    for evidence in section.evidences
                ),
            )
            for section in document.sections
        ),
    )


def _to_ppt_section(section: ReportSectionDraft, *, max_evidences: int) -> PptSection:
    return PptSection(
        category=section.category.value,
        title=section.title,
        status=section.status.value,
        importance_score=section.importance_score,
        impact_direction=section.impact_direction.value if section.impact_direction else None,
        time_horizon=section.time_horizon.value if section.time_horizon else None,
        current_summary=section.current_summary,
        key_facts=tuple(section.key_facts),
        historical_context=tuple(section.historical_context),
        implications=tuple(section.implications),
        watch_points=tuple(section.watch_points),
        evidences=tuple(_to_ppt_evidence(citation) for citation in section.news_citations[:max_evidences]),
    )


def _to_ppt_evidence(citation: ReportCitationDraft) -> PptEvidenceLine:
    return PptEvidenceLine(
        document_version_id=citation.document_version_id,
        citation_order=citation.citation_order,
        evidence_text=(citation.evidence_text or citation.document_title or "\ucd9c\ucc98 \uc815\ubcf4 \uc5c6\uc74c").strip(),
        relevance_score=citation.relevance_score,
        document_title=citation.document_title,
        source_name=citation.source_name,
        published_at=citation.published_at,
        source_url=citation.source_url,
    )


def _add_cover_slide(presentation: Presentation, document: PptReportDocument) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _fill_background(slide, BRAND_NAVY)
    _add_accent_band(slide, 0.0, 0.0, 0.34, document.layout.slide_height_inches, BRAND_ORANGE)
    _add_accent_band(slide, 0.34, 0.0, 0.12, document.layout.slide_height_inches, BRAND_GOLD)
    _add_logo(slide, left=SAFE_LEFT + 0.15, top=0.95, height=0.82)

    _add_text_block(
        slide,
        left=SAFE_LEFT + 0.15,
        top=2.35,
        width=SAFE_WIDTH_4_3 - 1.0,
        height=0.7,
        text=document.title,
        font_size=document.layout.title_font_size,
        bold=True,
        color=RGBColor(255, 255, 255),
        layout=document.layout,
    )
    _add_text_block(
        slide,
        left=SAFE_LEFT + 0.17,
        top=3.15,
        width=3.8,
        height=0.34,
        text="DAILY INDUSTRY BRIEF",
        font_size=12,
        bold=True,
        color=RGBColor(214, 222, 230),
        layout=document.layout,
    )
    _add_text_block(
        slide,
        left=SAFE_LEFT + 0.15,
        top=4.05,
        width=3.4,
        height=0.45,
        text=_format_header_date_dot(document.report_date),
        font_size=22,
        bold=True,
        color=RGBColor(255, 255, 255),
        layout=document.layout,
    )


def _add_agenda_slides(presentation: Presentation, document: PptReportDocument) -> None:
    if not document.sections:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        _fill_background(slide, BRAND_MIST)
        content_top = _add_slide_title(slide, "Agenda", document.layout)
        _add_text_block(
            slide,
            left=SAFE_LEFT,
            top=content_top + 0.35,
            width=SAFE_WIDTH_4_3,
            height=0.8,
            text="\uc644\ub8cc\ub41c \uc139\uc158\uc774 \uc5c6\uc2b5\ub2c8\ub2e4.",
            font_size=18,
            layout=document.layout,
            color=BRAND_TEXT,
        )
        return

    agenda_lines = [
        f"{index}. {_short_title(section.title, max_chars=34)}"
        for index, section in enumerate(document.sections, start=1)
    ]
    for chunk_index, items in enumerate(_chunked_text(agenda_lines, document.layout.max_agenda_items_per_slide), start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        _fill_background(slide, RGBColor(255, 252, 247))
        _add_accent_band(slide, 0.0, 0.0, 0.16, document.layout.slide_height_inches, BRAND_ORANGE)
        _add_logo(slide, left=8.42, top=0.52, height=0.34)
        title = "Agenda" if chunk_index == 1 else f"Agenda {chunk_index}"
        content_top = _add_slide_title(slide, title, document.layout)
        _add_bullet_card(slide, left=SAFE_LEFT, top=content_top + 0.25, width=4.25, height=4.85, title="\uc774\uc288 \ubaa9\ub85d", items=items, layout=document.layout)
        _add_bullet_card(
            slide,
            left=5.2,
            top=content_top + 0.25,
            width=4.25,
            height=4.85,
            title="\uc77d\ub294 \uc21c\uc11c",
            items=["\uc139\uc158 \uc694\uc57d\uc5d0\uc11c \uc774\uc288\ub97c \ud30c\uc545", "Highlights\uc5d0\uc11c \uc601\ud5a5 \ud655\uc778", "\ucd9c\ucc98 \ud45c\uc5d0\uc11c \uadfc\uac70 \ud655\uc778"],
            layout=document.layout,
        )


def _add_section_overview_slide(
    presentation: Presentation,
    document: PptReportDocument,
    index: int,
    section: PptSection,
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _fill_background(slide, RGBColor(255, 255, 255))
    _add_accent_band(slide, 0.0, 0.0, 0.12, document.layout.slide_height_inches, BRAND_GOLD)
    _add_logo(slide, left=8.42, top=0.52, height=0.34)
    content_top = _add_slide_title(slide, f"{index}. {_short_title(section.title)}", document.layout)
    _add_meta_chips(slide, section, document.layout, top=content_top + 0.18)
    _add_summary_panel(slide, section.current_summary or "\uc694\uc57d \uc815\ubcf4 \uc5c6\uc74c", document.layout, top=content_top + 0.78)
    _add_bullet_card(
        slide,
        left=6.55,
        top=content_top + 0.78,
        width=2.9,
        height=4.45,
        title="\ud575\uc2ec \uc2e0\ud638",
        items=[
            f"\uc911\uc694\ub3c4 {section.importance_score}" if section.importance_score is not None else "\uc911\uc694\ub3c4 n/a",
            f"\uc601\ud5a5 {section.impact_direction}" if section.impact_direction else "\uc601\ud5a5 n/a",
            f"\uc2dc\uac04 {section.time_horizon}" if section.time_horizon else "\uc2dc\uac04 n/a",
        ],
        layout=document.layout,
    )


def _add_section_highlight_slide(
    presentation: Presentation,
    document: PptReportDocument,
    index: int,
    section: PptSection,
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _fill_background(slide, BRAND_MIST)
    content_top = _add_slide_title(slide, f"{index}. {_short_title(section.title, max_chars=24)} | Highlights", document.layout)

    columns = [
        ("\ud575\uc2ec \uc0ac\uc2e4", _limit_bullets(section.key_facts)),
        ("SK\ud558\uc774\ub2c9\uc2a4 \uc601\ud5a5", _limit_bullets(section.implications)),
        ("\ub2e4\uc74c \ud655\uc778 \uc0ac\ud56d", _limit_bullets(section.watch_points or section.historical_context)),
    ]
    positions = [(SAFE_LEFT, content_top + 0.35), (3.65, content_top + 0.35), (6.75, content_top + 0.35)]
    for (title, items), (left, top) in zip(columns, positions):
        _add_bullet_card(slide, left=left, top=top, width=2.7, height=4.85, title=title, items=items or ["\ud45c\uc2dc\ud560 \ud56d\ubaa9 \uc5c6\uc74c"], layout=document.layout)


def _add_sources_slide(
    presentation: Presentation,
    document: PptReportDocument,
    chunk_index: int,
    source_chunk: list[PptEvidenceLine],
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _fill_background(slide, RGBColor(255, 255, 255))
    title = "\uc804\uccb4 \ucd9c\ucc98 \ubaa9\ub85d" if chunk_index == 1 else f"\uc804\uccb4 \ucd9c\ucc98 \ubaa9\ub85d {chunk_index}"
    content_top = _add_slide_title(slide, title, document.layout)
    _add_source_table(slide, source_chunk, document.layout, top=content_top + 0.28)


def _collect_source_items(sections: tuple[PptSection, ...]) -> list[PptEvidenceLine]:
    items: list[PptEvidenceLine] = []
    seen: set[tuple[str, str, str, str]] = set()
    for section in sections:
        for evidence in section.evidences:
            key = (evidence.source_name or "", evidence.document_title or evidence.evidence_text or "", evidence.published_at or "", evidence.source_url or "")
            if key in seen:
                continue
            seen.add(key)
            items.append(evidence)
    return items


def _format_source_line(section_title: str, evidence: PptEvidenceLine, source_index: int) -> str:
    title = _clean_internal_id(evidence.document_title or evidence.evidence_text or "\ucd9c\ucc98 \uc815\ubcf4 \uc5c6\uc74c")
    source = _clean_internal_id(evidence.source_name or "\ub274\uc2a4 \ucd9c\ucc98")
    published = _format_source_date(evidence.published_at)
    return f"{source_index}. {source} | {published} | {title}"


def _resolve_logo_path() -> Path | None:
    env_value = os.environ.get(PPT_LOGO_ENV)
    if env_value:
        candidate = Path(env_value)
        if candidate.exists():
            return candidate
    if DEFAULT_PPT_LOGO_PATH.exists():
        return DEFAULT_PPT_LOGO_PATH
    return None


def _add_logo(slide, *, left: float, top: float, height: float) -> None:
    logo_path = _resolve_logo_path()
    if logo_path is None:
        return
    slide.shapes.add_picture(str(logo_path), Inches(left), Inches(top), height=Inches(height))


def _add_accent_band(slide, left: float, top: float, width: float, height: float, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_info_card(slide, *, left: float, top: float, width: float, height: float, title: str, body: str, fill: RGBColor | None = None, title_color: RGBColor | None = None, body_color: RGBColor | None = None) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill or RGBColor(255, 255, 255)
    shape.line.color.rgb = RGBColor(227, 219, 207)
    frame = shape.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p1 = frame.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = DEFAULT_PPT_FONT_NAME
    r1.font.size = Pt(10)
    r1.font.bold = True
    r1.font.color.rgb = title_color or BRAND_MUTED
    p2 = frame.add_paragraph()
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = DEFAULT_PPT_FONT_NAME
    r2.font.size = Pt(16)
    r2.font.bold = True
    r2.font.color.rgb = body_color or BRAND_NAVY


def _add_meta_chips(slide, section: PptSection, layout: PptLayout, *, top: float) -> None:
    items = [f"\uce74\ud14c\uace0\ub9ac {section.category}"]
    if section.importance_score is not None:
        items.append(f"\uc911\uc694\ub3c4 {section.importance_score}")
    if section.impact_direction:
        items.append(f"\uc601\ud5a5 {section.impact_direction}")
    if section.time_horizon:
        items.append(f"\uc2dc\uac04 {section.time_horizon}")
    left = SAFE_LEFT
    for item in items:
        width = max(1.05, min(2.0, 0.1 * len(item) + 0.45))
        if left + width > 10.0 - SAFE_RIGHT:
            break
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.34))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(245, 247, 250)
        shape.line.color.rgb = RGBColor(221, 227, 233)
        frame = shape.text_frame
        frame.clear()
        frame.margin_left = Inches(0.08)
        frame.margin_right = Inches(0.08)
        p = frame.paragraphs[0]
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = _limit_text(item, 24)
        r.font.name = DEFAULT_PPT_FONT_NAME
        r.font.size = Pt(layout.meta_font_size)
        r.font.color.rgb = BRAND_MUTED
        left += width + 0.12


def _add_summary_panel(slide, summary_text: str, layout: PptLayout, *, top: float) -> None:
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(SAFE_LEFT), Inches(top), Inches(5.65), Inches(4.45))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    panel.line.color.rgb = RGBColor(225, 231, 236)
    frame = panel.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.18)
    frame.margin_right = Inches(0.18)
    frame.margin_top = Inches(0.16)
    p1 = frame.paragraphs[0]
    p1.space_after = Pt(6)
    r1 = p1.add_run()
    r1.text = "\ud604\uc7ac \uc694\uc57d"
    r1.font.name = DEFAULT_PPT_FONT_NAME
    r1.font.size = Pt(layout.card_title_font_size)
    r1.font.bold = True
    r1.font.color.rgb = BRAND_ORANGE
    p2 = frame.add_paragraph()
    p2.line_spacing = 1.18
    p2.space_after = Pt(4)
    r2 = p2.add_run()
    r2.text = _limit_text(summary_text, MAX_SUMMARY_CHARS)
    r2.font.name = DEFAULT_PPT_FONT_NAME
    r2.font.size = Pt(layout.body_font_size)
    r2.font.color.rgb = BRAND_TEXT


def _add_bullet_card(slide, *, left: float, top: float, width: float, height: float, title: str, items: list[str], layout: PptLayout) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = RGBColor(225, 231, 236)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.14)
    frame.margin_right = Inches(0.14)
    frame.margin_top = Inches(0.14)
    p1 = frame.paragraphs[0]
    p1.space_after = Pt(6)
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = DEFAULT_PPT_FONT_NAME
    r1.font.size = Pt(layout.card_title_font_size)
    r1.font.bold = True
    r1.font.color.rgb = BRAND_NAVY
    for item in _limit_bullets(items):
        p = frame.add_paragraph()
        p.text = item
        p.bullet = True
        p.level = 0
        p.line_spacing = 1.18
        p.space_after = Pt(4)
        p.font.name = DEFAULT_PPT_FONT_NAME
        p.font.size = Pt(layout.body_font_size)
        p.font.color.rgb = BRAND_TEXT


def _add_source_table(slide, sources: list[PptEvidenceLine], layout: PptLayout, *, top: float) -> None:
    rows = len(sources) + 1
    table_shape = slide.shapes.add_table(rows, 4, Inches(SAFE_LEFT), Inches(top), Inches(SAFE_WIDTH_4_3), Inches(0.42 * rows)).table
    headers = ["No", "\ucd9c\ucc98", "\uc81c\ubaa9", "\ub0a0\uc9dc"]
    widths = [0.48, 1.35, 5.55, 1.52]
    for index, width in enumerate(widths):
        table_shape.columns[index].width = Inches(width)
    for col_index, header in enumerate(headers):
        _set_table_cell(table_shape.cell(0, col_index), header, layout, bold=True, fill=BRAND_NAVY, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
    for row_index, source in enumerate(sources, start=1):
        _set_table_cell(table_shape.cell(row_index, 0), str(row_index), layout, align=PP_ALIGN.CENTER)
        _set_table_cell(table_shape.cell(row_index, 1), _limit_text(_clean_internal_id(source.source_name or "\ub274\uc2a4"), 20), layout)
        title = _clean_internal_id(source.document_title or source.evidence_text or "\ucd9c\ucc98 \uc815\ubcf4 \uc5c6\uc74c")
        _set_table_cell(table_shape.cell(row_index, 2), _limit_text(title, 60), layout)
        _set_table_cell(table_shape.cell(row_index, 3), _format_source_date(source.published_at), layout, align=PP_ALIGN.CENTER)


def _set_table_cell(cell, text: str, layout: PptLayout, *, bold: bool = False, fill: RGBColor | None = None, color: RGBColor | None = None, align: PP_ALIGN = PP_ALIGN.LEFT) -> None:
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    frame = cell.text_frame
    frame.clear()
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.04)
    p = frame.paragraphs[0]
    p.alignment = align
    p.line_spacing = 1.15
    p.space_after = Pt(0)
    r = p.add_run()
    r.text = text
    r.font.name = DEFAULT_PPT_FONT_NAME
    r.font.size = Pt(layout.meta_font_size)
    r.font.bold = bold
    r.font.color.rgb = color or BRAND_TEXT


def _fill_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_slide_title(slide, text: str, layout: PptLayout) -> float:
    title_text = _short_title(text, max_chars=46)
    title_height = 0.48 if len(title_text) <= 27 else TITLE_MAX_HEIGHT
    _add_text_block(
        slide,
        left=SAFE_LEFT,
        top=SAFE_TOP,
        width=SAFE_WIDTH_4_3,
        height=title_height,
        text=title_text,
        font_size=layout.slide_title_font_size,
        bold=True,
        color=BRAND_NAVY,
        layout=layout,
    )
    divider_top = SAFE_TOP + title_height + TITLE_DIVIDER_GAP
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(SAFE_LEFT),
        Inches(divider_top),
        Inches(SAFE_WIDTH_4_3),
        Inches(0.03),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BRAND_ORANGE
    line.line.fill.background()
    return divider_top + 0.03


def _add_text_block(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size: int | float,
    layout: PptLayout,
    bold: bool = False,
    color: RGBColor | None = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    word_wrap: bool = True,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = word_wrap
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.name = DEFAULT_PPT_FONT_NAME
    font.size = Pt(font_size)
    font.bold = bold
    if color is not None:
        font.color.rgb = color


def _add_bullet_block(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    items: list[str],
    font_size: int | float,
    layout: PptLayout,
    color: RGBColor | None = None,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.bullet = True
        paragraph.font.name = DEFAULT_PPT_FONT_NAME
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color or BRAND_TEXT


def _short_title(value: str, *, max_chars: int = 38) -> str:
    cleaned = " ".join((value or "").split())
    if len(cleaned) <= max_chars:
        return cleaned
    cut = cleaned[:max_chars].rstrip(" ,.;:\u00b7\u318d-/")
    if " " in cut:
        cut = cut.rsplit(" ", 1)[0]
    return f"{cut}..."


def _limit_text(value: str, max_chars: int) -> str:
    cleaned = _clean_internal_id(" ".join((value or "").split()))
    if len(cleaned) <= max_chars:
        return cleaned
    return f"{cleaned[: max_chars - 3].rstrip()}..."


def _limit_bullets(items: tuple[str, ...] | list[str]) -> list[str]:
    return [_limit_text(item, MAX_BULLET_CHARS) for item in list(items)[:MAX_COLUMN_ITEMS] if str(item).strip()]


def _clean_internal_id(value: str) -> str:
    cleaned = re.sub(r"[0-9a-fA-F]{8}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}-[0-9a-fA-F]{12}", "", value or "")
    cleaned = re.sub(r"\b(?:daily|weekly|report|workspace|analysis|doc|ver|version|format|pptx)[:/_-][\w:./-]+", "", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"\b(?:completed|Citation reference|Format|Version|UUID)\b", "", cleaned, flags=re.IGNORECASE)
    return " ".join(cleaned.replace("[]", "").split()).strip()


def _format_source_date(value: str | None) -> str:
    if not value:
        return "-"
    return value[:10].replace("-", ".")


def _format_header_date_dot(value: str) -> str:
    parsed = _parse_date_text(value)
    if parsed is None:
        return value.replace("-", ".")
    return f"{parsed.year:04d}.{parsed.month:02d}.{parsed.day:02d}"


def _format_header_date_korean(value: str) -> str:
    parsed = _parse_date_text(value)
    if parsed is None:
        return value
    return f"{parsed.year}\ub144 {parsed.month}\uc6d4 {parsed.day}\uc77c"


def _parse_date_text(value: str) -> date | None:
    raw = (value or "").strip()
    if not raw:
        return None
    candidate = raw.split("T", 1)[0].replace(".", "-")
    try:
        return date.fromisoformat(candidate)
    except ValueError:
        return None


def _chunked(values: tuple[PptEvidenceLine, ...], chunk_size: int) -> list[tuple[PptEvidenceLine, ...]]:
    return [values[index : index + chunk_size] for index in range(0, len(values), chunk_size)]


def _chunked_text(values: list[str], chunk_size: int) -> list[list[str]]:
    return [values[index : index + chunk_size] for index in range(0, len(values), chunk_size)]


def _format_evidence_line(evidence: PptEvidenceLine) -> str:
    base = _clean_internal_id(evidence.document_title or evidence.evidence_text or "\ucd9c\ucc98 \uc815\ubcf4 \uc5c6\uc74c")
    if evidence.relevance_score is None:
        return base
    return f"{base} (\uad00\ub828\ub3c4: {evidence.relevance_score:.2f})"
