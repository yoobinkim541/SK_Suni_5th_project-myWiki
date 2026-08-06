from __future__ import annotations

from datetime import date
from io import BytesIO

from docx import Document
from pptx import Presentation
from pypdf import PdfReader
import pytest

from src.report.interface import ReportGenerationConfig, generate_daily_report
from src.report.models import (
    ArtifactType,
    GeneratedReport,
    ReportGenerationRequest,
    ReportSectionDraft,
    ReportStatus,
)
from src.analysis.importance_models import ImpactDirection, TimeHorizon
from src.analysis.models import Category
from src.report.repository import SavedReportArtifact


def make_section() -> ReportSectionDraft:
    return ReportSectionDraft(
        issue_key="issue-default-artifacts",
        representative_analysis_result_id="analysis-1",
        category=Category.PRODUCT_TECHNOLOGY,
        importance_score=90,
        impact_direction=ImpactDirection.RISK,
        time_horizon=TimeHorizon.MID_TERM,
        title="샘플 섹션 제목",
        current_summary="샘플 보고서 본문",
        status="completed",
    )


def make_artifact(artifact_type: ArtifactType) -> SavedReportArtifact:
    extension = {ArtifactType.PDF: "pdf", ArtifactType.DOCX: "docx", ArtifactType.PPTX: "pptx"}[artifact_type]
    return SavedReportArtifact(
        artifact_id=f"artifact-{artifact_type.value}",
        report_id="report-1",
        artifact_type=artifact_type,
        object_key=f"ws-1/reports/report-1/{artifact_type.value}/v1.{extension}",
        version=1,
        mime_type="application/octet-stream",
        file_size=123,
    )


def test_report_generation_config_defaults_to_downloadable_artifacts() -> None:
    config = ReportGenerationConfig()

    assert config.artifacts.formats == [ArtifactType.PDF, ArtifactType.DOCX, ArtifactType.PPTX]


def test_report_artifact_config_none_uses_downloadable_artifacts() -> None:
    config = ReportGenerationConfig(artifacts={"formats": None})

    assert config.artifacts.formats == [ArtifactType.PDF, ArtifactType.DOCX, ArtifactType.PPTX]


def test_generate_daily_report_creates_pdf_docx_and_pptx_by_default(monkeypatch: pytest.MonkeyPatch) -> None:
    created_artifact_types: list[ArtifactType] = []

    monkeypatch.setattr("src.report.interface.create_report_version", lambda **kwargs: GeneratedReport(
        report_id="report-1",
        workspace_id="ws-1",
        report_date=date(2026, 8, 3),
        report_type="daily",
        version=1,
        status=ReportStatus.PENDING,
        created_at="2026-08-03T08:00:00+00:00",
    ))
    monkeypatch.setattr("src.report.interface.get_report_candidates", lambda **kwargs: [object()])
    monkeypatch.setattr("src.report.interface.select_report_candidates", lambda *args, **kwargs: [object()])
    monkeypatch.setattr("src.report.interface.group_report_candidates", lambda *args, **kwargs: [object()])
    monkeypatch.setattr("src.report.interface.enrich_issue_groups", lambda *args, **kwargs: [object()])
    monkeypatch.setattr("src.report.interface.compose_report_sections", lambda *args, **kwargs: [make_section()])
    monkeypatch.setattr("src.report.interface.generate_wiki_drafts_for_sections", lambda *args, **kwargs: [])
    monkeypatch.setattr("src.report.interface.assemble_generated_report", lambda **kwargs: GeneratedReport(
        report_id="report-1",
        workspace_id="ws-1",
        report_date=date(2026, 8, 3),
        report_type="daily",
        status=ReportStatus.DRAFTING,
        title="일일 보고서",
        language="ko",
        sections=[make_section()],
    ))
    monkeypatch.setattr("src.report.interface.save_report_sections", lambda **kwargs: [object()])
    monkeypatch.setattr("src.report.interface.mark_report_completed", lambda **kwargs: None)
    monkeypatch.setattr("src.report.interface.create_and_save_markdown_artifact", lambda **kwargs: (_ for _ in ()).throw(RuntimeError("should not be called")))

    def capture_pdf(**kwargs):
        created_artifact_types.append(ArtifactType.PDF)
        return make_artifact(ArtifactType.PDF)

    def capture_docx(**kwargs):
        created_artifact_types.append(ArtifactType.DOCX)
        return make_artifact(ArtifactType.DOCX)

    def capture_pptx(**kwargs):
        created_artifact_types.append(ArtifactType.PPTX)
        return make_artifact(ArtifactType.PPTX)

    monkeypatch.setattr("src.report.interface.create_and_save_pdf_artifact", capture_pdf)
    monkeypatch.setattr("src.report.interface.create_and_save_docx_artifact", capture_docx)
    monkeypatch.setattr("src.report.interface.create_and_save_pptx_artifact", capture_pptx)

    result = generate_daily_report(
        ReportGenerationRequest(workspace_id="ws-1", report_date=date(2026, 8, 3), max_sections=5),
    )

    assert created_artifact_types == [ArtifactType.PDF, ArtifactType.DOCX, ArtifactType.PPTX]
    assert [artifact.artifact_type for artifact in result.artifacts] == created_artifact_types
    assert result.artifact.artifact_type == ArtifactType.PDF
    assert result.report.artifact_type == ArtifactType.PDF

class FakeResult:
    def __init__(self, data):
        self.data = data


class FakeTable:
    def __init__(self, supabase: "FakeSupabase", name: str):
        self.supabase = supabase
        self.name = name
        self.rows = supabase.tables.setdefault(name, [])
        self.eq_filters: list[tuple[str, object]] = []
        self.insert_rows: list[dict[str, object]] | None = None

    def select(self, _fields: str) -> "FakeTable":
        return self

    def eq(self, field: str, value: object) -> "FakeTable":
        self.eq_filters.append((field, value))
        return self

    def insert(self, rows):
        if isinstance(rows, dict):
            rows = [rows]
        self.insert_rows = [dict(row) for row in rows]
        return self

    def execute(self) -> FakeResult:
        if self.insert_rows is not None:
            inserted = []
            for row in self.insert_rows:
                new_row = dict(row)
                new_row.setdefault("id", f"{self.name}-{len(self.rows) + 1}")
                self.rows.append(new_row)
                inserted.append(dict(new_row))
            return FakeResult(inserted)

        rows = self.rows
        for field, value in self.eq_filters:
            rows = [row for row in rows if row.get(field) == value]
        return FakeResult([dict(row) for row in rows])


class FakeStorageBucket:
    def __init__(self, supabase: "FakeSupabase", bucket_name: str):
        self.supabase = supabase
        self.bucket_name = bucket_name

    def upload(self, *, path: str, file: bytes, file_options: dict[str, object] | None = None):
        bucket_objects = self.supabase.storage_objects.setdefault(self.bucket_name, {})
        bucket_objects[path] = file
        return {"path": path}

    def download(self, path: str) -> bytes:
        return self.supabase.storage_objects[self.bucket_name][path]


class FakeStorage:
    def __init__(self, supabase: "FakeSupabase"):
        self.supabase = supabase

    def from_(self, bucket_name: str) -> FakeStorageBucket:
        self.supabase.storage_objects.setdefault(bucket_name, {})
        return FakeStorageBucket(self.supabase, bucket_name)


class FakeSupabase:
    def __init__(self):
        self.tables: dict[str, list[dict[str, object]]] = {"artifacts": []}
        self.storage = FakeStorage(self)
        self.storage_objects: dict[str, dict[str, bytes]] = {}

    def table(self, name: str) -> FakeTable:
        return FakeTable(self, name)


def extract_pdf_text(payload: bytes) -> str:
    reader = PdfReader(BytesIO(payload))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def extract_docx_text(payload: bytes) -> str:
    document = Document(BytesIO(payload))
    return "\n".join(paragraph.text for paragraph in document.paragraphs)


def extract_pptx_text(payload: bytes) -> str:
    presentation = Presentation(BytesIO(payload))
    return "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )


def test_generate_daily_report_default_artifacts_include_report_content(monkeypatch: pytest.MonkeyPatch) -> None:
    supabase = FakeSupabase()
    section = make_section()

    monkeypatch.setattr("src.report.interface.create_report_version", lambda **kwargs: GeneratedReport(
        report_id="report-1",
        workspace_id="ws-1",
        report_date=date(2026, 8, 3),
        report_type="daily",
        version=1,
        status=ReportStatus.PENDING,
        created_at="2026-08-03T08:00:00+00:00",
    ))
    monkeypatch.setattr("src.report.interface.get_report_candidates", lambda **kwargs: [object()])
    monkeypatch.setattr("src.report.interface.select_report_candidates", lambda *args, **kwargs: [object()])
    monkeypatch.setattr("src.report.interface.group_report_candidates", lambda *args, **kwargs: [object()])
    monkeypatch.setattr("src.report.interface.enrich_issue_groups", lambda *args, **kwargs: [object()])
    monkeypatch.setattr("src.report.interface.compose_report_sections", lambda *args, **kwargs: [section])
    monkeypatch.setattr("src.report.interface.generate_wiki_drafts_for_sections", lambda *args, **kwargs: [])
    monkeypatch.setattr("src.report.interface.assemble_generated_report", lambda **kwargs: GeneratedReport(
        workspace_id="ws-1",
        report_date=date(2026, 8, 3),
        report_type="daily",
        status=ReportStatus.DRAFTING,
        title="일일 산업 동향 보고서",
        language="ko",
        sections=[section],
    ))
    monkeypatch.setattr("src.report.interface.save_report_sections", lambda **kwargs: [object()])
    monkeypatch.setattr("src.report.interface.mark_report_completed", lambda **kwargs: None)
    monkeypatch.setattr("src.report.interface.create_and_save_markdown_artifact", lambda **kwargs: (_ for _ in ()).throw(RuntimeError("should not be called")))

    result = generate_daily_report(
        ReportGenerationRequest(workspace_id="ws-1", report_date=date(2026, 8, 3), max_sections=5),
        supabase=supabase,
    )

    assert [artifact.artifact_type for artifact in result.artifacts] == [ArtifactType.PDF, ArtifactType.DOCX, ArtifactType.PPTX]

    payloads = {
        artifact.artifact_type: supabase.storage_objects["reports"][artifact.object_key]
        for artifact in result.artifacts
    }
    extracted = {
        ArtifactType.PDF: extract_pdf_text(payloads[ArtifactType.PDF]),
        ArtifactType.DOCX: extract_docx_text(payloads[ArtifactType.DOCX]),
        ArtifactType.PPTX: extract_pptx_text(payloads[ArtifactType.PPTX]),
    }

    for artifact_type, text in extracted.items():
        assert "일일 산업 동향 보고서" in text, artifact_type
        assert "샘플 섹션 제목" in text, artifact_type
        assert "샘플 보고서 본문" in text, artifact_type
