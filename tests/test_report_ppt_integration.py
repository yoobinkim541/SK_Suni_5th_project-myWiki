from __future__ import annotations

from copy import deepcopy
from datetime import date, datetime, timezone
from io import BytesIO

from pptx import Presentation
import pytest

from src.analysis.importance_models import ImpactDirection, TimeHorizon
from src.analysis.models import Category
from src.report.artifact_service import (
    DEFAULT_REPORT_ARTIFACT_BUCKET,
    PPTX_CONTENT_TYPE,
    create_and_save_pptx_artifact,
    save_pptx_report_artifact,
)
from src.report.interface import ReportGenerationConfig, generate_daily_report
from src.report.models import (
    ArtifactType,
    GeneratedReport,
    ReportCitationDraft,
    ReportGenerationRequest,
    ReportSectionDraft,
    ReportStatus,
    ReportType,
)
from src.report.repository import SavedReportArtifact


class FakeResult:
    def __init__(self, data):
        self.data = data


class FakeTable:
    def __init__(self, supabase: "FakeSupabase", name: str):
        self.supabase = supabase
        self.name = name
        self.rows = supabase.tables.setdefault(name, [])
        self.eq_filters: list[tuple[str, object]] = []
        self.insert_rows: list[dict[str, object]] | None = None

    def select(self, _fields: str) -> "FakeTable":
        return self

    def eq(self, field: str, value: object) -> "FakeTable":
        self.eq_filters.append((field, value))
        return self

    def insert(self, rows):
        if isinstance(rows, dict):
            rows = [rows]
        self.insert_rows = [dict(row) for row in rows]
        return self

    def execute(self) -> FakeResult:
        if self.insert_rows is not None:
            inserted = []
            for row in self.insert_rows:
                new_row = dict(row)
                new_row.setdefault("id", f"{self.name}-{len(self.rows) + 1}")
                self.rows.append(new_row)
                inserted.append(dict(new_row))
            return FakeResult(inserted)
        rows = self.rows
        for field, value in self.eq_filters:
            rows = [row for row in rows if row.get(field) == value]
        return FakeResult([dict(row) for row in rows])


class FakeStorageBucket:
    def __init__(self, supabase: "FakeSupabase", bucket_name: str):
        self.supabase = supabase
        self.bucket_name = bucket_name

    def upload(self, *, path: str, file: bytes, file_options: dict[str, object] | None = None):
        self.supabase.last_upload = {
            "bucket": self.bucket_name,
            "path": path,
            "file": file,
            "file_options": dict(file_options or {}),
        }
        bucket_objects = self.supabase.storage_objects.setdefault(self.bucket_name, {})
        bucket_objects[path] = file
        return {"path": path}

    def download(self, path: str) -> bytes:
        return self.supabase.storage_objects[self.bucket_name][path]

    def remove(self, paths: list[str]):
        bucket_objects = self.supabase.storage_objects.setdefault(self.bucket_name, {})
        for path in paths:
            bucket_objects.pop(path, None)
        return {"paths": paths}


class FakeStorage:
    def __init__(self, supabase: "FakeSupabase"):
        self.supabase = supabase

    def from_(self, bucket_name: str) -> FakeStorageBucket:
        self.supabase.storage_objects.setdefault(bucket_name, {})
        return FakeStorageBucket(self.supabase, bucket_name)


class FakeSupabase:
    def __init__(self):
        self.tables: dict[str, list[dict[str, object]]] = {"artifacts": []}
        self.storage = FakeStorage(self)
        self.storage_objects: dict[str, dict[str, bytes]] = {}
        self.last_upload: dict[str, object] | None = None

    def table(self, name: str) -> FakeTable:
        return FakeTable(self, name)


def make_section() -> ReportSectionDraft:
    return ReportSectionDraft(
        issue_key="issue-1",
        representative_analysis_result_id="analysis-1",
        category=Category.PRODUCT_TECHNOLOGY,
        importance_score=88,
        impact_direction=ImpactDirection.OPPORTUNITY,
        time_horizon=TimeHorizon.MID_TERM,
        title="\uc0d8\ud50c \uc139\uc158 \uc81c\ubaa9",
        current_summary="\uc0d8\ud50c \ubcf8\ubb38\uc785\ub2c8\ub2e4.",
        key_facts=["\ud575\uc2ec \uc0ac\uc2e4"],
        implications=["\uc2dc\uc0ac\uc810"],
        watch_points=["\uad00\ucc30 \ud3ec\uc778\ud2b8"],
        news_citations=[
            ReportCitationDraft(
                analysis_result_id="analysis-1",
                document_version_id="doc-ver-1",
                citation_order=1,
                evidence_text="\uadfc\uac70 \ubb38\uc7a5",
            )
        ],
        status="completed",
    )


def make_report() -> GeneratedReport:
    return GeneratedReport(
        report_id="report-1",
        workspace_id="ws-001",
        report_date=date(2026, 8, 3),
        report_type=ReportType.DAILY,
        title="\uc77c\uc77c \ubcf4\uace0\uc11c",
        language="ko",
        version=3,
        status=ReportStatus.COMPLETED,
        sections=[make_section()],
        created_at=datetime(2026, 8, 3, 8, 15, tzinfo=timezone.utc),
        generated_at=datetime(2026, 8, 3, 8, 15, tzinfo=timezone.utc),
    )


def make_artifact(*, artifact_type: str = "pptx", object_key: str = "ws-1/reports/report-1/pptx/v1.pptx") -> SavedReportArtifact:
    return SavedReportArtifact(
        artifact_id=f"artifact-{artifact_type}",
        report_id="report-1",
        artifact_type=artifact_type,
        object_key=object_key,
        version=1,
        mime_type="application/octet-stream",
        file_size=123,
        content_hash="abc123",
        storage_bucket="reports",
    )


def test_save_pptx_report_artifact_uploads_and_persists_metadata() -> None:
    supabase = FakeSupabase()
    payload = b"PK\x03\x04pptx-bytes"

    artifact = save_pptx_report_artifact(
        report_id="report-1",
        workspace_id="ws-001",
        report_version=3,
        report_date=date(2026, 8, 3),
        pptx_bytes=payload,
        supabase=supabase,
    )

    assert artifact.object_key == "ws-001/reports/report-1/pptx/v3.pptx"
    assert artifact.mime_type == PPTX_CONTENT_TYPE
    assert artifact.storage_bucket == DEFAULT_REPORT_ARTIFACT_BUCKET
    assert supabase.last_upload is not None
    assert supabase.last_upload["path"] == "ws-001/reports/report-1/pptx/v3.pptx"
    assert supabase.last_upload["file"] == payload


def test_create_and_save_pptx_artifact_uses_renderer_when_bytes_missing() -> None:
    supabase = FakeSupabase()
    report = make_report()

    artifact = create_and_save_pptx_artifact(report=report, supabase=supabase)

    assert artifact.object_key.endswith(".pptx")
    uploaded = supabase.last_upload
    assert uploaded is not None
    presentation = Presentation(BytesIO(uploaded["file"]))
    text = "\n".join(shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert "\uc77c\uc77c \uc0b0\uc5c5 \ub3d9\ud5a5 \ubcf4\uace0\uc11c" in text
    assert "2026\ub144 8\uc6d4 3\uc77c" in text
    assert "\uc0d8\ud50c \uc139\uc158 \uc81c\ubaa9" in text


def test_create_and_save_pptx_artifact_does_not_mutate_inputs() -> None:
    supabase = FakeSupabase()
    report = make_report()
    original = deepcopy(report.model_dump())

    create_and_save_pptx_artifact(report=report, supabase=supabase)

    assert report.model_dump() == original


def test_generate_daily_report_can_create_pptx_artifact(monkeypatch: pytest.MonkeyPatch) -> None:
    captured: dict[str, object] = {}
    config = ReportGenerationConfig(artifacts={"formats": [ArtifactType.PPTX]})

    monkeypatch.setattr("src.report.interface.create_report_version", lambda **kwargs: GeneratedReport(
        report_id="report-1",
        workspace_id="ws-1",
        report_date=date(2026, 8, 3),
        report_type="daily",
        version=1,
        status=ReportStatus.PENDING,
        created_at="2026-08-03T08:00:00+00:00",
    ))
    monkeypatch.setattr("src.report.interface.get_report_candidates", lambda **kwargs: [object()])
    monkeypatch.setattr("src.report.interface.select_report_candidates", lambda *args, **kwargs: [object()])
    monkeypatch.setattr("src.report.interface.group_report_candidates", lambda *args, **kwargs: [object()])
    monkeypatch.setattr("src.report.interface.enrich_issue_groups", lambda *args, **kwargs: [object()])
    monkeypatch.setattr("src.report.interface.compose_report_sections", lambda *args, **kwargs: [make_section()])
    monkeypatch.setattr("src.report.interface.generate_wiki_drafts_for_sections", lambda *args, **kwargs: [])
    monkeypatch.setattr("src.report.interface.assemble_generated_report", lambda **kwargs: GeneratedReport(
        report_id="report-1",
        workspace_id="ws-1",
        report_date=date(2026, 8, 3),
        report_type="daily",
        status=ReportStatus.DRAFTING,
        title="\uc77c\uc77c \ubcf4\uace0\uc11c",
        language="ko",
        sections=[make_section()],
    ))
    monkeypatch.setattr("src.report.interface.save_report_sections", lambda **kwargs: [object()])
    monkeypatch.setattr(
        "src.report.interface.create_and_save_pptx_artifact",
        lambda **kwargs: captured.update(kwargs) or make_artifact(),
    )
    monkeypatch.setattr("src.report.interface.create_and_save_markdown_artifact", lambda **kwargs: (_ for _ in ()).throw(RuntimeError("should not be called")))
    monkeypatch.setattr("src.report.interface.create_and_save_docx_artifact", lambda **kwargs: (_ for _ in ()).throw(RuntimeError("should not be called")))
    monkeypatch.setattr("src.report.interface.mark_report_completed", lambda **kwargs: None)

    result = generate_daily_report(
        ReportGenerationRequest(workspace_id="ws-1", report_date=date(2026, 8, 3), max_sections=5),
        config=config,
    )

    assert captured["created_by"] is None
    assert result.artifact.artifact_type == ArtifactType.PPTX
    assert result.artifacts[0].artifact_type == ArtifactType.PPTX
    assert result.report.artifact_type == ArtifactType.PPTX
