from __future__ import annotations

from datetime import date
from io import BytesIO

from pptx import Presentation

from src.analysis.importance_models import ImpactDirection, TimeHorizon
from src.analysis.models import Category
from src.report.models import ReportCitationDraft, ReportSectionDraft, ReportSectionStatus
from src.report.ppt_renderer import (
    DEFAULT_PPT_TITLE,
    build_daily_report_ppt_document,
    build_daily_report_ppt_filename,
    render_daily_report_ppt,
)

K_TITLE = "\uc77c\uc77c \uc0b0\uc5c5 \ub3d9\ud5a5 \ub9de\ucda4 \ubcf4\uace0\uc11c"
K_SECTION = "HBM \uc218\uc694 \ud655\ub300\uc640 \uacf5\uae09\ub9dd \ud22c\uc790 \uac00\uc18d"
K_BODY = "\uad6d\ub0b4 \uba54\ubaa8\ub9ac \uc81c\uc870\uc0ac\uc640 \uacf5\uae09\ub9dd \uae30\uc5c5\ub4e4\uc774 AI \uc11c\ubc84 \uc218\uc694 \uc99d\uac00\uc5d0 \ub300\uc751\ud558\uae30 \uc704\ud574 \uc0dd\uc0b0\ub2a5\ub825 \ud655\uc7a5\uacfc \uc7a5\ube44 \ud22c\uc790\ub97c \ub3d9\uc2dc\uc5d0 \ucd94\uc9c4\ud558\uace0 \uc788\uc2b5\ub2c8\ub2e4."
K_EVIDENCE = "HBM \uc218\uc694 \ud655\ub300\uc5d0 \ub9de\ucdb0 \uad00\ub828 \uacf5\uae09\ub9dd \uc804\ubc18\uc5d0\uc11c \ud22c\uc790 \uc18d\ub3c4\uac00 \ube68\ub77c\uc9c0\uace0 \uc788\ub2e4."


def make_section(
    *,
    status: ReportSectionStatus = ReportSectionStatus.COMPLETED,
    title: str = K_SECTION,
) -> ReportSectionDraft:
    return ReportSectionDraft(
        issue_key="issue-1",
        representative_analysis_result_id="analysis-1",
        category=Category.PRODUCT_TECHNOLOGY,
        importance_score=88,
        impact_direction=ImpactDirection.OPPORTUNITY,
        time_horizon=TimeHorizon.MID_TERM,
        title=title,
        current_summary=K_BODY,
        key_facts=["\ud575\uc2ec \uc0ac\uc2e4 1", "\ud575\uc2ec \uc0ac\uc2e4 2"],
        implications=["SK\ud558\uc774\ub2c9\uc2a4 \uc218\ud61c\uc0ac \ud655\ub300"],
        watch_points=["\uad00\ucc30 \ud3ec\uc778\ud2b8"],
        news_citations=[
            ReportCitationDraft(
                analysis_result_id="analysis-1",
                document_version_id="doc-ver-1",
                citation_order=1,
                evidence_text=K_EVIDENCE,
                relevance_score=0.92,
                source_name="\ub274\uc2a4\uc6d0",
                published_at="2026-08-03T09:00:00+09:00",
            )
        ],
        status=status,
    )


def collect_ppt_text(pptx_bytes: bytes) -> str:
    presentation = Presentation(BytesIO(pptx_bytes))
    texts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            texts.append(cell_text)
            elif hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    texts.append(text)
    return "\n".join(texts)


def test_build_daily_report_ppt_document_uses_completed_sections_only() -> None:
    completed = make_section(status=ReportSectionStatus.COMPLETED)
    pending = make_section(status=ReportSectionStatus.DRAFTING)

    document = build_daily_report_ppt_document(
        report_key="daily-trends-2026-08-03",
        version=2,
        sections=[completed, pending],
        generated_at="2026-08-03T09:00:00+09:00",
        report_date=date(2026, 8, 3),
        title=K_TITLE,
    )

    assert document.title == K_TITLE
    assert document.subtitle == "daily-trends-2026-08-03"
    assert document.report_date == "2026-08-03"
    assert document.version == 2
    assert len(document.sections) == 1
    assert document.sections[0].title == K_SECTION


def test_build_daily_report_ppt_filename_normalizes_spaces() -> None:
    assert build_daily_report_ppt_filename(report_key="daily trends 2026-08-03", version=3) == "daily-trends-2026-08-03-v3.pptx"


def test_render_daily_report_ppt_returns_pptx_bytes() -> None:
    report_document = build_daily_report_ppt_document(
        report_key="daily-trends-2026-08-03",
        version=1,
        sections=[make_section()],
        generated_at="2026-08-03T09:00:00+09:00",
        report_date=date(2026, 8, 3),
        title=K_TITLE,
    )

    pptx_bytes = render_daily_report_ppt(report_document)
    text = collect_ppt_text(pptx_bytes)

    assert pptx_bytes[:2] == b"PK"
    assert K_TITLE in text
    assert "Mywiki" not in text
    assert "Version" not in text
    assert "Format" not in text
    assert "daily-trends-2026-08-03" not in text
    assert "2026.08.03" in text
    assert K_SECTION in text
    assert K_BODY in text
    assert K_EVIDENCE in text


def test_render_daily_report_ppt_prefers_report_date_over_generated_at() -> None:
    report_document = build_daily_report_ppt_document(
        report_key="daily-trends-2026-08-03",
        version=1,
        sections=[make_section()],
        generated_at="2026-08-02T23:30:00+00:00",
        report_date=date(2026, 8, 3),
    )

    pptx_bytes = render_daily_report_ppt(report_document)
    text = collect_ppt_text(pptx_bytes)

    assert "2026.08.03" in text


def test_render_daily_report_ppt_creates_agenda_and_evidence_slides() -> None:
    report_document = build_daily_report_ppt_document(
        report_key="daily-trends-2026-08-03",
        version=1,
        sections=[make_section()],
        report_date=date(2026, 8, 3),
    )

    presentation = Presentation(BytesIO(render_daily_report_ppt(report_document)))

    assert len(presentation.slides) >= 4
    assert presentation.slide_width == 9144000
    assert presentation.slide_height == 6858000


def test_render_daily_report_ppt_keeps_wrapped_highlight_title_above_divider() -> None:
    long_title = "SK\ud558\uc774\ub2c9\uc2a4 \ubbf8\uad6d \uc0c1\uc7a5 \ud6c4 \uc6d4\uac00 \uae0d\uc815 \uc804\ub9dd\uacfc AI \uba54\ubaa8\ub9ac \uc218\uc694 \uad6c\uc870\uc801 \ubcc0\ud654 \ubc0f \ucc28\uc138\ub300 \uc2dc\uc7a5 \uacbd\uc7c1 \uc2ec\ud654"
    report_document = build_daily_report_ppt_document(
        report_key="daily-trends-2026-08-03",
        version=1,
        sections=[make_section(title=long_title)],
        report_date=date(2026, 8, 3),
    )

    presentation = Presentation(BytesIO(render_daily_report_ppt(report_document)))
    highlight_slide = next(
        slide
        for slide in presentation.slides
        if any(getattr(shape, "text", "").startswith("1.") and "Highlights" in getattr(shape, "text", "") for shape in slide.shapes)
    )
    shapes = list(highlight_slide.shapes)
    title_index = next(index for index, shape in enumerate(shapes) if getattr(shape, "text", "").startswith("1.") and "Highlights" in getattr(shape, "text", ""))
    title_shape = shapes[title_index]
    divider_shape = shapes[title_index + 1]

    assert long_title not in title_shape.text
    assert title_shape.text.count("\n") <= 1
    assert title_shape.height <= 731520
    assert title_shape.top + title_shape.height <= divider_shape.top
